#!/usr/bin/env python
# coding: utf-8

import pptx
import pandas as pd
import os
import sys

class PptEditor2:

    #Output a list of paragraph texts
    def pptx_to_text_list(self, pptx_path):
        prs = pptx.Presentation(pptx_path)
        text_list = []
        for slide_i, slide in enumerate(prs.slides):
                for shape_i, shape in enumerate(slide.shapes):
                    if not shape.has_text_frame:
                        continue
                    for para_i, para in enumerate(shape.text_frame.paragraphs):
                        text_list.append((f"Slide {slide_i+1}",para.text))
        return text_list


    def compare_pptx(self, pptx_paths):

        pptx_list = []
        for path in pptx_paths:
            pptx_list.append(self.pptx_to_text_list(path))

        df = pd.DataFrame()
        df["Slide No"] = [pair[0] for pair in pptx_list[0]]
        for index, text_list in enumerate(pptx_list):
            df[f"File {str(index+1)}: \n{os.path.basename(pptx_paths[index])}"] = [pair[1] for pair in text_list]

            
        directory, file_name = os.path.split(pptx_paths[0])
        file, ext = os.path.splitext(file_name)

        df.to_csv(os.path.join(directory, f"{file}_compared.csv"), encoding="utf-8-sig", index=False)

        print(df)
        
        return df

if __name__ == '__main__':
    editor2 = PptEditor2()
    editor2.compare_pptx((sys.argv[1:]))