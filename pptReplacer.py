#!/usr/bin/env python
# coding: utf-8

import pptx
import os
import pandas as pd
import datetime
import sys


class pptReplacer:

    def __init__(self, ppt_path, csv_path):
        self.prs = pptx.Presentation(ppt_path)
        self.dict_df = pd.read_csv(csv_path)

        self.directory, self.file_name = os.path.split(ppt_path)
        self.file, self.ext = os.path.splitext(self.file_name)

    def replaceWord(self, before_after_cols=["Run Text", "Revised Run Text"]):  # FAILED???

        # fldr_path, file_name = os.path.dirname(file_path), os.path.basename(file_path)
        # prs = pptx.Presentation(file_path)

        # dict_df = pd.read_csv(dict_file)
        # dict_df[before_after_cols[0]] = dict_df[before_after_cols[0]].astype(str)
        # dict_df[before_after_cols[1]] = dict_df[before_after_cols[1]].astype(str)
        # dict_df[before_after_cols[1]] = dict_df[before_after_cols[1]].fillna("")
        # dictionary = dict(zip(dict_df[before_after_cols[0]], dict_df[before_after_cols[1]]))

        # for ns, slide in enumerate(prs.slides):
        #     for nsh, shape in enumerate(slide.shapes):
        #         if not shape.has_text_frame:
        #             continue
        #         for np, paragraph in enumerate(shape.text_frame.paragraphs):
        #             for rs, run in enumerate(paragraph.runs):
        #                 passage_before = run.text
        #
        #                 #                     #Replace a holw run text by comparing "訳文" and "修正"
        #                 #                     for key in dictionary.keys():
        #                 #                         if key in passage_before and key != dictionary.get(key):
        #                 #                             passage_after = passage_before.replace(key, dictionary.get(key))
        #                 #                             prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = passage_after
        #                 #                             print("passage_before: ", passage_before)
        #                 #                             print("passage_after: ", passage_after)
        #                 #                             print("-" * 100)
        #                 #                             passage_before = passage_after
        #                 for key in dictionary.keys():
        #                     if key == passage_before != dictionary.get(key):
        #                         passage_after = dictionary.get(key)
        #                         prs.slides[ns].shapes[nsh].text_frame.paragraphs[np].runs[rs].text = passage_after
        #                         print("passage_before: ", passage_before)
        #                         print("passage_after: ", passage_after)
        #                         print("-" * 100)

        before_col = self.dict_df[before_after_cols[0]]
        after_col = self.dict_df[before_after_cols[1]]
        for index, element in enumerate(before_col):
            if element != after_col[index]:
                slide_i = self.dict_df["Slide No"][index]
                shape_i = self.dict_df["Shape No"][index]
                para_i = self.dict_df["Paragraph No"][index]
                run_i = self.dict_df["Run No"][index]

                before_text = self.prs.slides[slide_i].shapes[shape_i].text_frame.paragraphs[para_i].runs[run_i].text
                self.prs.slides[slide_i].shapes[shape_i].text_frame.paragraphs[para_i].runs[run_i].text = after_col[index]
                after_text = self.prs.slides[slide_i].shapes[shape_i].text_frame.paragraphs[para_i].runs[run_i].text
                print(f"before_text: {before_text}")
                print(f"after_text : {after_text}")
                print("-"*100)


        # Get the latest timestamp for the new file name
        timestamp = datetime.datetime.now()
        timestamp = str(timestamp)[2:16]
        table = str.maketrans({'-': '',
                               ' ': '-',
                               ':': '',
                               })
        timestamp = timestamp.translate(table)

        # self.prs.save(self.directory + "\\" + self.file_name[:-5] + "_" + timestamp + ".pptx")
        self.prs.save(self.directory + "\\" + self.file + ".pptx")





if __name__ == '__main__':
    editor = pptReplacer(sys.argv[1], sys.argv[2])
    editor.replaceWord()


    # print(sys.argv[1:])
    # for i in sys.argv[1:]:
    #     print(i)