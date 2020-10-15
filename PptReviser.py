#!/usr/bin/env python
# coding: utf-8

import pptx
import pandas as pd
import os
import sys

class PptReviser:

    def __init__(self, pptx_path):
        self.prs = pptx.Presentation(pptx_path)
        self.directory, self.file_name = os.path.split(pptx_path)
        self.file, self.ext = os.path.splitext(self.file_name)

    #Output a list of run texts
    def pptx_to_text_list(self, pptx_path):
        text_list = []
        for slide_i, slide in enumerate(self.prs.slides):
                for shape_i, shape in enumerate(slide.shapes):
                    if not shape.has_text_frame:
                        continue
                    for para_i, para in enumerate(shape.text_frame.paragraphs):
                        for run_i, run in enumerate(para.runs):
                            text_list.append((slide_i,
                                              shape_i,
                                              para_i,
                                              para.text,
                                              run_i,
                                              run.text))
        return text_list


    def list_run_text(self, pptx_path, pptx_path2):
        prs2 = pptx.Presentation(pptx_path2)

        text_list = self.pptx_to_text_list(pptx_path)

        df = pd.DataFrame()
        cols = ["Slide No", "Shape No", "Paragraph No", "Paragraph Text", "Run No", "Run Text"]
        for col_i, col in enumerate(cols):
            df[col] = [pair[col_i] for pair in text_list]
        df[f"Revised {col}"] = df[col] #Add "Revised Run Text" Column

        #Insert original pptx's paragraph texts
        df.insert(3, "Original Paragraph Text", "orig para text")
        for i in range(df.shape[0]):
            slide_i = df["Slide No"][i]
            shape_i = df["Shape No"][i]
            para_i = df["Paragraph No"][i]
            # df.iloc[i, 3] = prs2.slides[slide_i].shapes[shape_i].text_frame.paragraphs[para_i].text
            df["Original Paragraph Text"][i] = prs2.slides[slide_i].shapes[shape_i].text_frame.paragraphs[para_i].text

        df.to_csv(os.path.join(self.directory, f"{self.file}_revised.csv"), encoding="utf-8-sig", index=False)

        print(df)

        return df


if __name__ == '__main__':
    editor = PptReviser((sys.argv[1]))
    editor.list_run_text(sys.argv[1], sys.argv[2])